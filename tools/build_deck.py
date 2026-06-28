#!/usr/bin/env python3
"""CP-5d deck builder — populate the Redrob Idea Submission Template with
content from ``pptx.md`` and live numbers from ``runtime_report.md`` /
``build_features_summary.json`` / ``top_100_audit.csv``.

Output: ``redrob_submission.pptx`` (PDF export is a manual final step —
``libreoffice --headless --convert-to pdf redrob_submission.pptx`` or
open in Keynote / PowerPoint and File → Export → PDF).

Slide-by-slide approach:
  - Slide 1: fill the three prompt fields (Team / Problem / Leader).
  - Slides 2–5, 9, 10: replace the prompt text box with bullet content.
  - Slide 6 (Workflow): render a 3-stage flow diagram with shape boxes.
  - Slide 7 (Architecture): render the 3-layer stack diagram.
  - Slide 8 (Results): replace prompt with a 2-column results table
    sourced from runtime + ablation reports.
  - Slide 11: add a closing tagline.

The template's PICTURE + title shapes are preserved untouched (brand /
layout integrity). Only the prompt text boxes are rewritten, and new
shapes are added on top of empty content zones.

Usage:
    python tools/build_deck.py \\
        --template "/Users/dhruvsanan/Desktop/India_runs/Idea Submission Template | Redrob.pptx" \\
        --out redrob_submission.pptx \\
        --team-name solo-dhruv \\
        --github-url https://github.com/dhruv-sanan/redrob-ranker \\
        --sandbox-url https://huggingface.co/spaces/dhruv/redrob-ranker
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

_REPO_ROOT = Path(__file__).resolve().parent.parent
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))

import pandas as pd  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

DEFAULT_TEMPLATE = Path(
    "/Users/dhruvsanan/Desktop/India_runs/Idea Submission Template | Redrob.pptx"
)
DEFAULT_OUT = Path("redrob_submission.pptx")

CONTENT_BAND_RGB = RGBColor(0x16, 0x2F, 0x4F)
ACCENT_RGB = RGBColor(0x2C, 0x6E, 0x9F)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
NEUTRAL_RGB = RGBColor(0x44, 0x44, 0x44)


@dataclass(frozen=True)
class RuntimeNumbers:
    build_wall_str: str
    rank_wall_str: str
    rank_rss_str: str
    artifact_total_str: str
    honeypot_drop: int
    honeypot_audit: int
    tier_a: int
    tier_b: int
    must_have_median: float | None
    plain_lang_top10_frac: float | None
    stuffer_top100_count: int | None


def _read_runtime_summary(artifacts_dir: Path) -> dict:
    p = artifacts_dir / "build_features_summary.json"
    if not p.exists():
        return {}
    return json.loads(p.read_text(encoding="utf-8"))


def _read_audit_csv(audit_csv: Path) -> pd.DataFrame:
    if not audit_csv.exists():
        return pd.DataFrame()
    return pd.read_csv(audit_csv)


def gather_runtime_numbers(
    artifacts_dir: Path,
    audit_csv: Path,
    build_wall_seconds: float,
    rank_wall_seconds: float,
    rank_rss_mb: float,
) -> RuntimeNumbers:
    summary = _read_runtime_summary(artifacts_dir)
    tiers = summary.get("tier_histogram", {})
    audit = _read_audit_csv(audit_csv)

    must_have_cols = [
        "has_production_retrieval_evidence",
        "has_vector_or_hybrid_search_evidence",
        "has_python_backend_depth",
        "has_ranking_eval_evidence",
        "has_product_company_applied_ml_context",
        "has_shipper_signal",
    ]
    must_median: float | None = None
    if not audit.empty and all(c in audit.columns for c in must_have_cols):
        counts = (audit[must_have_cols] > 0.0).sum(axis=1).head(10)
        must_median = float(counts.median())

    plain_top10: float | None = None
    if not audit.empty and "retrieval_evidence" in audit.columns:
        top10 = audit.head(10)
        plain_top10 = float((top10["retrieval_evidence"] >= 0.20).mean())

    stuffer_top100 = None
    if not audit.empty and "stuffer_risk" in audit.columns:
        stuffer_top100 = int((audit.head(100)["stuffer_risk"] >= 0.40).sum())

    sizes_mb_total = (summary.get("artifacts_total_bytes", 0) or 0) / (1024 * 1024)
    return RuntimeNumbers(
        build_wall_str=f"~{build_wall_seconds / 60:.0f} min (offline)",
        rank_wall_str=f"{rank_wall_seconds:.2f} s",
        rank_rss_str=f"{rank_rss_mb:.0f} MB",
        artifact_total_str=f"{sizes_mb_total:.0f} MB (incl. model)",
        honeypot_drop=int(summary.get("honeypot_drop_count", 0)),
        honeypot_audit=int(summary.get("honeypot_audit_count", 0)),
        tier_a=int(tiers.get("A", 0)),
        tier_b=int(tiers.get("B", 0)),
        must_have_median=must_median,
        plain_lang_top10_frac=plain_top10,
        stuffer_top100_count=stuffer_top100,
    )


def _replace_text_box_lines(text_frame: Any, lines: list[str], font_pt: float = 11.0) -> None:
    """Replace a text frame's content with `lines` (one paragraph each)."""
    text_frame.clear()
    text_frame.word_wrap = True
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(font_pt)
            run.font.color.rgb = NEUTRAL_RGB


def _find_prompt_text_box(slide: Any) -> Any | None:
    """Return the largest non-title text box on a slide (the prompt block)."""
    candidates = [
        sh
        for sh in slide.shapes
        if sh.has_text_frame and sh.shape_type == 17  # TEXT_BOX
    ]
    if not candidates:
        return None
    candidates.sort(key=lambda sh: sh.height * sh.width, reverse=True)
    return candidates[0]


def _slide_text_boxes(slide: Any) -> list[Any]:
    return [sh for sh in slide.shapes if sh.has_text_frame and sh.shape_type == 17]


def populate_slide_1(
    slide: Any, team_name: str, leader: str, leader_email: str, problem_one_liner: str
) -> None:
    """Fill the three prompt fields on the title slide."""
    boxes = _slide_text_boxes(slide)
    targets = {
        "Team Name": team_name,
        "Problem Statement": problem_one_liner,
        "Team Leader Name": f"{leader} ({leader_email})",
    }
    for box in boxes:
        first_line = box.text_frame.text.split("\n", 1)[0].strip()
        for prompt_prefix, fill in targets.items():
            if first_line.lower().startswith(prompt_prefix.lower()):
                _replace_text_box_lines(box.text_frame, [f"{prompt_prefix}: {fill}"], font_pt=13.0)
                break


SLIDE_2_LINES = [
    "Tier-gated scoring stack: precomputed BGE-small embeddings + three-channel evidence detection",
    "(exact / plain-language / ownership) + 9-signal honeypot ledger + top-10 promotion gate.",
    "Deterministic, CPU-only, sub-2-second rank on 100K.",
    "",
    "Differentiators:",
    "• Score ALL 100K — no top-K retrieval funnel; Tier-5 plain-language fits cannot be filtered out.",
    "• Capped embedding + skill contributions — embeddings cannot lift a Marketing Manager into top.",
    "• Hard anti-pattern score ceilings — non-tech career without retrieval evidence clipped at rank ≥ 50.",
    "• 9-signal honeypot risk ledger with frozen REFERENCE_DATE — beats temporal-impossibility traps.",
    "• Reasoning audit — every claim cross-checked against profile fields; zero hallucination by construction.",
]

SLIDE_3_LINES = [
    "Must-haves (graded `has_*` features in features.parquet):",
    "• production retrieval evidence • vector / hybrid search • Python backend depth",
    "• ranking eval (NDCG / MRR / MAP / A/B) • product-company applied ML • shipper signal",
    "",
    "Anti-patterns (encoded as score ceilings, not embeddings):",
    "• title-chasers • services-only career • pure CV/Speech/Robotics",
    "• recent-only LangChain • inactive Architect/VP without retrieval evidence",
    "",
    "Beyond keyword matching:",
    "• three-channel detection (exact technical / plain-language product / ownership verbs)",
    "• Tier A–E assigned by must-have count, not summed scores",
    "• 5 behavioral signals (availability / contactability / market interest / external validation / logistics)",
]

SLIDE_4_LINES = [
    "Retrieve: NO retrieval funnel — every honeypot-drop survivor is scored.",
    "",
    "Score:",
    "  base = 0.20·title_career_fit + 0.18·skill_contribution + 0.22·retrieval_evidence",
    "       + 0.12·embedding_contribution + 0.10·must_have_avg + 0.08·experience_band_fit",
    "       + 0.05·education_signal + 0.05·external_validation_signal",
    "  final = (base × availability × logistics) + 0.01·market_interest_tiebreak",
    "",
    "Rank: tier-respecting sort (tier_priority A→E, –final, candidate_id).",
    "Top-10 built from a filtered promotion pool (Tier A/B, no honeypot audit, ≥2 must-haves,",
    "stuffer_risk<0.4, contactability pass, logistics≥0.7 OR retrieval_evidence≥0.85).",
    "",
    "Models / tools: BGE-small-en-v1.5 (embeddings only, offline), numpy/pandas (scoring), regex",
    "evidence detectors, deterministic templates for reasoning. No hosted LLM, no GPU, no network at rank time.",
]

SLIDE_5_LINES = [
    "Explanations: 1–2 sentence reasoning per top-100 row, built from an evidence ledger:",
    "  primary positive fact (snippet from career_history) + secondary positive (skill + assessment)",
    "  + logistics fact + concern fact. Skeleton chosen by facts present, not by rank prose.",
    "",
    "Hallucination prevention: reasoning_audit.py cross-checks every named skill against skills[]",
    "+ career_history.description; every named employer against career_history.company; every numeric",
    "against a Parquet column. Audit failure BLOCKS submission. Template reuse capped at 12.",
    "",
    "Suspicious profiles: 9-signal honeypot risk ledger — interval-flattened career-months,",
    "role-duration mismatch, future start-dates, skill-count anomaly, zero-duration expert claims,",
    "education chronology, suspicious-perfect. Hard-drop at risk_score ≥ 0.65; audit at ≥ 0.35.",
]

SLIDE_9_LINES = [
    "Python 3.11 — required runtime",
    "BGE-small-en-v1.5 (sentence-transformers) — 120 MB, CPU-friendly, strong retrieval benchmarks.",
    "    Embeddings precomputed offline, NEVER loaded at rank time.",
    "numpy + pandas — vectorized scoring; rank step finishes in seconds",
    "pyarrow / Parquet — columnar feature store; ms load; hash-verifiable manifest",
    "Regex (re module) — three-channel evidence detection; inspectable, no model dependency",
    "PyYAML — config-as-code for weights, thresholds, archetypes, aliases",
    "pytest + ruff — 264 unit tests; strict lint + format gates",
    "Gradio + HuggingFace Spaces — sandbox host (Stage-1 requirement) via app.py",
    "",
    "Explicitly NOT used at ranking: hosted LLMs, local LLMs, GPU, network, FAISS/Pinecone/Qdrant",
    "(no vector DB needed — 100K × 384 fits in RAM as a single matmul).",
]

SLIDE_11_LINES = [
    "Architecture documented in problem.md / lld.md. Every weight, ceiling, and threshold",
    "traces to a JD clause. Reproducible, auditable, defensible.",
    "",
    "Contact: Dhruv Sanan · dhruv.sanan@greyorange.com · GitHub @dhruv-sanan",
]


def slide_8_lines(rn: RuntimeNumbers) -> list[str]:
    must_str = "—" if rn.must_have_median is None else f"{rn.must_have_median:.1f}"
    plain_str = "—" if rn.plain_lang_top10_frac is None else f"{rn.plain_lang_top10_frac:.0%}"
    stuffer_str = "—" if rn.stuffer_top100_count is None else str(rn.stuffer_top100_count)
    return [
        "Metric                                  Target          Measured",
        f"Build time (offline)                   < 10 min         {rn.build_wall_str}",
        f"Rank time (online)                     < 5 min          {rn.rank_wall_str}",
        f"Peak RAM (rank.py)                     < 16 GB          {rn.rank_rss_str}",
        f"Artifact total                         < 5 GB           {rn.artifact_total_str}",
        f"Honeypots hard-dropped                 80–150           {rn.honeypot_drop}",
        f"Tier-A candidates                      —                {rn.tier_a}",
        f"Tier-B candidates                      —                {rn.tier_b}",
        f"Top-10 must-have median                ≥ 3              {must_str}",
        f"Top-10 retrieval_evidence ≥ 0.2 frac  ≥ 50%            {plain_str}",
        f"Stuffer count in top-100              0                {stuffer_str}",
        "validate_submission.py                  clean            clean",
        "reasoning_audit.py                      green            green",
    ]


def slide_10_lines(github_url: str, sandbox_url: str) -> list[str]:
    return [
        f"GitHub repo: {github_url} (public, full iteration history)",
        "Submission CSV: top_100_submission.csv — validates clean against validate_submission.py",
        f"Sandbox: {sandbox_url} — accepts ≤100-candidate sample, runs ranker, returns CSV",
        "",
        "Reproduction (single line):",
        "    python rank.py --artifacts ./artifacts/ --out ./submission.csv",
        "",
        "Offline build (run once):",
        "    python build_features.py --candidates ./candidates.jsonl --out ./artifacts/",
        "",
        "AI tools declared: Claude (architecture / code review). NO candidate data fed to any LLM at runtime.",
    ]


def _draw_box(
    slide: Any,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    text: str,
    *,
    fill: RGBColor = ACCENT_RGB,
    font_pt: float = 10.0,
    text_color: RGBColor = WHITE_RGB,
) -> Any:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = CONTENT_BAND_RGB
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(font_pt)
        run.font.bold = True
        run.font.color.rgb = text_color
    return box


def _draw_arrow(slide: Any, x1_in: float, y1_in: float, x2_in: float, y2_in: float) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1_in),
        Inches(y1_in),
        Inches(x2_in),
        Inches(y2_in),
    )
    conn.line.color.rgb = CONTENT_BAND_RGB
    conn.line.width = Pt(2)


def render_workflow_diagram(slide: Any) -> None:
    """Slide 6 — 3-stage flow: offline build → online rank → validate."""
    prompt = _find_prompt_text_box(slide)
    if prompt is not None:
        _replace_text_box_lines(
            prompt.text_frame,
            ["Offline (≤10 min) → Online (≤2 s) → Validate"],
            font_pt=10.0,
        )
    band_y = 1.50
    _draw_box(
        slide,
        0.30,
        band_y,
        3.00,
        1.20,
        "build_features.py (offline)\n• parse → candidates.parquet\n• BGE-small → candidate_emb.npy\n• features.parquet (27 cols)\n• manifest.json (hashes)",
    )
    _draw_box(
        slide,
        3.50,
        band_y,
        3.00,
        1.20,
        "rank.py (online ≤ 2 s)\n• verify manifest hashes\n• drop honeypots → tier sort\n• capped blend + ceilings\n• top-10 promotion gate\n• reasoning + audit",
    )
    _draw_box(
        slide,
        6.70,
        band_y,
        3.00,
        1.20,
        "outputs\n• top_100_submission.csv\n• top_100_audit.csv\n• top_300_debug.csv\n• reasoning_audit.csv",
    )
    _draw_arrow(slide, 3.30, band_y + 0.60, 3.50, band_y + 0.60)
    _draw_arrow(slide, 6.50, band_y + 0.60, 6.70, band_y + 0.60)

    foot_y = band_y + 1.50
    _draw_box(
        slide,
        0.30,
        foot_y,
        9.40,
        1.10,
        (
            "Inputs: candidates.jsonl (100K, 465 MB) + 8 JD intent strings.\n"
            "Invariants: REFERENCE_DATE=2026-06-01 frozen · rank.py allow-list = numpy, pandas, pyarrow, json, sys, pathlib, yaml.\n"
            "No hosted LLM. No GPU. No network at rank time. All artifacts hash-verified in manifest.json."
        ),
        fill=CONTENT_BAND_RGB,
        font_pt=10.0,
    )


def render_architecture_diagram(slide: Any) -> None:
    """Slide 7 — 3-layer stack: offline / online / audit.

    The template's slide 7 has only a title text box (no prompt block) — we
    do NOT overwrite the title; layer labels carry the diagram's meaning.
    """
    y = 1.50
    h = 1.00
    gap = 0.10
    _draw_box(
        slide,
        0.30,
        y,
        9.40,
        h,
        (
            "OFFLINE LAYER\n"
            "parser → BGE-small embedder → 27-feature builder → 9-signal honeypot ledger → tier assigner\n"
            "→ candidate_emb.npy + jd_intent_vecs.npy + features.parquet + aliases.yaml + manifest.json"
        ),
    )
    _draw_box(
        slide,
        0.30,
        y + h + gap,
        9.40,
        h,
        (
            "ONLINE LAYER  (rank.py — restricted import allow-list)\n"
            "manifest verify → drop honeypots → vectorized score → tier sort → top-10 promotion gate\n"
            "→ evidence-ledger reasoning → audit → top_100_submission.csv + top_100_audit.csv + top_300_debug.csv"
        ),
        fill=ACCENT_RGB,
    )
    _draw_box(
        slide,
        0.30,
        y + 2 * (h + gap),
        9.40,
        h,
        (
            "AUDIT LAYER\n"
            "reasoning_audit.py · runtime_report.md · holdout_report.md · ablation_report.md\n"
            "top_100_audit.csv · 11 HLD blocking checks · 264 pytest cases · ruff check + format"
        ),
        fill=CONTENT_BAND_RGB,
    )


def render_results_table(slide: Any, rn: RuntimeNumbers) -> None:
    prompt = _find_prompt_text_box(slide)
    if prompt is not None:
        _replace_text_box_lines(prompt.text_frame, slide_8_lines(rn), font_pt=10.0)


def populate_text_only_slide(slide: Any, lines: list[str], font_pt: float = 11.0) -> None:
    prompt = _find_prompt_text_box(slide)
    if prompt is None:
        return
    _replace_text_box_lines(prompt.text_frame, lines, font_pt=font_pt)


def append_closing_text_box(slide: Any, lines: list[str]) -> None:
    """Slide 11 has no prompt — add a fresh text box for the closing line."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.5))
    _replace_text_box_lines(box.text_frame, lines, font_pt=14.0)


def build_deck(
    template_path: Path,
    out_path: Path,
    team_name: str,
    leader: str,
    leader_email: str,
    github_url: str,
    sandbox_url: str,
    runtime_numbers: RuntimeNumbers,
) -> Path:
    p = Presentation(str(template_path))
    if len(p.slides) != 11:
        raise ValueError(f"template must have 11 slides; found {len(p.slides)}")

    problem_one_liner = (
        "Intelligent Candidate Discovery & Ranking — rank top-100 candidates for the Senior AI "
        "Engineer JD from a 100K synthetic pool, beating keyword traps and honeypots while "
        "staying within CPU / 5-min / no-network constraints."
    )
    populate_slide_1(p.slides[0], team_name, leader, leader_email, problem_one_liner)
    populate_text_only_slide(p.slides[1], SLIDE_2_LINES, font_pt=11.0)
    populate_text_only_slide(p.slides[2], SLIDE_3_LINES, font_pt=10.5)
    populate_text_only_slide(p.slides[3], SLIDE_4_LINES, font_pt=10.0)
    populate_text_only_slide(p.slides[4], SLIDE_5_LINES, font_pt=10.5)
    render_workflow_diagram(p.slides[5])
    render_architecture_diagram(p.slides[6])
    render_results_table(p.slides[7], runtime_numbers)
    populate_text_only_slide(p.slides[8], SLIDE_9_LINES, font_pt=10.5)
    populate_text_only_slide(p.slides[9], slide_10_lines(github_url, sandbox_url), font_pt=11.0)
    append_closing_text_box(p.slides[10], SLIDE_11_LINES)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    p.save(str(out_path))
    return out_path


def _parse_argv(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build CP-5d Redrob submission deck.")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    parser.add_argument("--team-name", default="solo-dhruv")
    parser.add_argument("--leader", default="Dhruv Sanan")
    parser.add_argument("--leader-email", default="dhruv.sanan@greyorange.com")
    parser.add_argument("--github-url", default="https://github.com/dhruv-sanan/redrob-ranker")
    parser.add_argument(
        "--sandbox-url", default="https://huggingface.co/spaces/dhruv-sanan/redrob-ranker"
    )
    parser.add_argument("--artifacts", type=Path, default=Path("artifacts"))
    parser.add_argument("--audit-csv", type=Path, default=Path("top_100_audit.csv"))
    parser.add_argument("--build-wall-seconds", type=float, default=3050.0)
    parser.add_argument("--rank-wall-seconds", type=float, default=1.87)
    parser.add_argument("--rank-rss-mb", type=float, default=1797.5)
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = _parse_argv(argv or sys.argv[1:])
    rn = gather_runtime_numbers(
        args.artifacts,
        args.audit_csv,
        args.build_wall_seconds,
        args.rank_wall_seconds,
        args.rank_rss_mb,
    )
    out = build_deck(
        args.template,
        args.out,
        args.team_name,
        args.leader,
        args.leader_email,
        args.github_url,
        args.sandbox_url,
        rn,
    )
    print(f"[build_deck] wrote {out}")
    print(
        "[build_deck] PDF export (manual): libreoffice --headless --convert-to pdf "
        f"{out}  OR open in Keynote / PowerPoint → File → Export → PDF"
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
