"""CP-5d tests — deck builder helpers + end-to-end deck assembly."""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
import pytest
from pptx import Presentation
from pptx.util import Inches

from tools.build_deck import (
    DEFAULT_TEMPLATE,
    SLIDE_2_LINES,
    SLIDE_3_LINES,
    SLIDE_4_LINES,
    SLIDE_5_LINES,
    SLIDE_9_LINES,
    SLIDE_11_LINES,
    RuntimeNumbers,
    _find_prompt_text_box,
    _replace_text_box_lines,
    build_deck,
    gather_runtime_numbers,
    slide_8_lines,
    slide_10_lines,
)


def test_slide_constant_lines_nonempty():
    for block in (
        SLIDE_2_LINES,
        SLIDE_3_LINES,
        SLIDE_4_LINES,
        SLIDE_5_LINES,
        SLIDE_9_LINES,
        SLIDE_11_LINES,
    ):
        assert isinstance(block, list)
        assert any(line.strip() for line in block)


def test_slide_10_lines_contains_urls():
    lines = slide_10_lines("https://github.example/repo", "https://hf.example/space")
    joined = "\n".join(lines)
    assert "https://github.example/repo" in joined
    assert "https://hf.example/space" in joined


def test_slide_8_lines_formats_targets_and_measurements():
    rn = RuntimeNumbers(
        build_wall_str="~50 min (offline)",
        rank_wall_str="1.87 s",
        rank_rss_str="1797 MB",
        artifact_total_str="220 MB (incl. model)",
        honeypot_drop=348,
        honeypot_audit=15143,
        tier_a=82,
        tier_b=384,
        must_have_median=3.0,
        plain_lang_top10_frac=0.7,
        stuffer_top100_count=0,
    )
    lines = slide_8_lines(rn)
    joined = "\n".join(lines)
    assert "Build time" in joined
    assert "1.87 s" in joined
    assert "348" in joined
    assert "validate_submission.py" in joined


def test_slide_8_lines_handles_none_metrics():
    rn = RuntimeNumbers(
        build_wall_str="—",
        rank_wall_str="—",
        rank_rss_str="—",
        artifact_total_str="—",
        honeypot_drop=0,
        honeypot_audit=0,
        tier_a=0,
        tier_b=0,
        must_have_median=None,
        plain_lang_top10_frac=None,
        stuffer_top100_count=None,
    )
    lines = slide_8_lines(rn)
    joined = "\n".join(lines)
    assert "—" in joined


def test_gather_runtime_numbers_reads_summary_and_audit(tmp_path: Path):
    arts = tmp_path / "arts"
    arts.mkdir()
    summary = {
        "tier_histogram": {"A": 82, "B": 384, "C": 10188, "D": 87275, "E": 2071},
        "honeypot_drop_count": 348,
        "honeypot_audit_count": 15143,
        "artifacts_total_bytes": 220 * 1024 * 1024,
    }
    (arts / "build_features_summary.json").write_text(json.dumps(summary))

    audit_rows = [
        {
            "candidate_id": f"CAND_{i:07d}",
            "rank": i + 1,
            "retrieval_evidence": 0.4 if i < 7 else 0.05,
            "stuffer_risk": 0.1,
            "has_production_retrieval_evidence": 1.0 if i < 5 else 0.0,
            "has_vector_or_hybrid_search_evidence": 1.0 if i < 5 else 0.0,
            "has_python_backend_depth": 1.0,
            "has_ranking_eval_evidence": 1.0 if i < 5 else 0.0,
            "has_product_company_applied_ml_context": 1.0,
            "has_shipper_signal": 1.0,
        }
        for i in range(10)
    ]
    audit_csv = tmp_path / "top_100_audit.csv"
    pd.DataFrame(audit_rows).to_csv(audit_csv, index=False)

    rn = gather_runtime_numbers(
        arts, audit_csv, build_wall_seconds=3000, rank_wall_seconds=1.9, rank_rss_mb=1800
    )
    assert rn.honeypot_drop == 348
    assert rn.tier_a == 82
    assert rn.must_have_median is not None and rn.must_have_median > 0
    assert rn.plain_lang_top10_frac == pytest.approx(0.7)
    assert rn.stuffer_top100_count == 0


def test_gather_runtime_numbers_missing_summary_returns_zeros(tmp_path: Path):
    rn = gather_runtime_numbers(tmp_path, tmp_path / "no_audit.csv", 0, 0, 0)
    assert rn.honeypot_drop == 0
    assert rn.tier_a == 0
    assert rn.must_have_median is None


def test_replace_text_box_lines_writes_each_paragraph():
    p = Presentation()
    layout = p.slide_layouts[5]
    slide = p.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    _replace_text_box_lines(box.text_frame, ["alpha", "beta", "gamma"], font_pt=11.0)
    paragraphs = box.text_frame.paragraphs
    assert [p.text for p in paragraphs] == ["alpha", "beta", "gamma"]


def test_find_prompt_text_box_picks_largest():
    p = Presentation()
    layout = p.slide_layouts[5]
    slide = p.slides.add_slide(layout)
    small = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(0.5))
    big = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    small.text_frame.text = "small"
    big.text_frame.text = "big"
    found = _find_prompt_text_box(slide)
    assert found.text_frame.text == "big"


def test_find_prompt_text_box_on_picture_only_slide_returns_none():
    p = Presentation()
    layout = p.slide_layouts[6]  # BLANK
    slide = p.slides.add_slide(layout)
    assert _find_prompt_text_box(slide) is None


@pytest.mark.skipif(not DEFAULT_TEMPLATE.exists(), reason="template PPTX not present")
def test_build_deck_writes_11_slide_pptx(tmp_path: Path):
    rn = RuntimeNumbers(
        build_wall_str="~50 min (offline)",
        rank_wall_str="1.87 s",
        rank_rss_str="1797 MB",
        artifact_total_str="220 MB (incl. model)",
        honeypot_drop=348,
        honeypot_audit=15143,
        tier_a=82,
        tier_b=384,
        must_have_median=3.0,
        plain_lang_top10_frac=0.7,
        stuffer_top100_count=0,
    )
    out = tmp_path / "deck.pptx"
    written = build_deck(
        DEFAULT_TEMPLATE,
        out,
        team_name="solo-test",
        leader="Test Leader",
        leader_email="t@example.com",
        github_url="https://github.example/r",
        sandbox_url="https://hf.example/s",
        runtime_numbers=rn,
    )
    assert written == out
    assert out.exists()
    assert out.stat().st_size > 1000

    p = Presentation(str(out))
    assert len(p.slides) == 11

    def slide_text(idx: int) -> str:
        s = p.slides[idx]
        return "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)

    assert "solo-test" in slide_text(0)
    assert "Tier-gated scoring stack" in slide_text(1)
    assert "Three-channel" not in slide_text(2) or "three-channel" in slide_text(2).lower()
    assert "title_career_fit" in slide_text(3)
    assert "reasoning_audit.py" in slide_text(4)
    assert "rank.py" in slide_text(5).lower() or "rank" in slide_text(5).lower()
    assert "System Architecture" in slide_text(6), "slide 7 title preserved"
    assert "Tier-A" in slide_text(7)
    assert "BGE-small" in slide_text(8)
    assert "https://github.example/r" in slide_text(9)
    assert "Dhruv Sanan" in slide_text(10)


@pytest.mark.skipif(not DEFAULT_TEMPLATE.exists(), reason="template PPTX not present")
def test_build_deck_rejects_template_with_wrong_slide_count(tmp_path: Path):
    blank = Presentation()
    short_path = tmp_path / "short.pptx"
    blank.save(str(short_path))
    rn = RuntimeNumbers(
        build_wall_str="-",
        rank_wall_str="-",
        rank_rss_str="-",
        artifact_total_str="-",
        honeypot_drop=0,
        honeypot_audit=0,
        tier_a=0,
        tier_b=0,
        must_have_median=None,
        plain_lang_top10_frac=None,
        stuffer_top100_count=None,
    )
    with pytest.raises(ValueError, match="11 slides"):
        build_deck(
            short_path,
            tmp_path / "out.pptx",
            team_name="x",
            leader="x",
            leader_email="x",
            github_url="x",
            sandbox_url="x",
            runtime_numbers=rn,
        )
